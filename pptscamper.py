import requests
from pptx import Presentation
from pptx.util import Inches

# Function to download images
def download_image(url, filename):
    response = requests.get(url)
    if response.status_code == 200:
        with open(filename, 'wb') as f:
            f.write(response.content)

# Create a presentation object
prs = Presentation()

# Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "SCAMPER: Innovating Cricket Bats"
subtitle.text = "A Creative Approach to Design\nYour Name / Date"
# Download and add image for title slide
download_image('https://example.com/cricket_bat.jpg', 'cricket_bat.jpg')
slide_1.shapes.add_picture('cricket_bat.jpg', Inches(3), Inches(1.5), width=Inches(4.5))

# Slide 1: Introduction to SCAMPER
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Introduction to SCAMPER"
content.text = ("Definition of SCAMPER: A creative thinking technique.\n"
                "Application in design thinking: Generates innovative ideas.\n"
                "Purpose: Innovate the design and functionality of cricket bats.")
# Download and add SCAMPER diagram image
download_image('https://www.creativefabrica.com/wp-content/uploads/2020/11/SCAMPER-Design-Thinking-Diagram-Graphic-By-CreativeFabrica-1.jpg', 'scamper_diagram.jpg')
slide_2.shapes.add_picture('scamper_diagram.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 2: Substitute
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "Substitute"
content.text = ("What materials can be substituted in the cricket bat?\n"
                "Example: Using carbon fiber instead of traditional willow or ash.\n"
                "Benefits: Increased durability and lighter weight.")
# Download and add material comparison image
download_image('https://www.shutterstock.com/image-photo/wooden-cricket-bat-isolated-on-white-600w-763373717.jpg', 'material_comparison.jpg')
slide_3.shapes.add_picture('material_comparison.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 3: Combine
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]
title.text = "Combine"
content.text = ("How can we combine features to enhance the cricket bat?\n"
                "Example: Traditional bat with technology such as sensors for swing analysis.\n"
                "Benefits: Improved performance analysis.")
# Download and add smart bat image
download_image('https://cdn.shopify.com/s/files/1/0020/2275/9246/products/2019-21BB-210-01.jpg', 'smart_cricket_bat.jpg')
slide_4.shapes.add_picture('smart_cricket_bat.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 4: Adapt
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Adapt"
content.text = ("What features from other sports can be adapted for cricket bats?\n"
                "Example: Grip technology from tennis rackets.\n"
                "Benefits: Better comfort and control.")
# Download and add grip image
download_image('https://img.freepik.com/free-photo/tennis-racket-background-with-wooden-bat_23-2148671392.jpg', 'grip_comparison.jpg')
slide_5.shapes.add_picture('grip_comparison.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 5: Modify (Magnify)
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]
title.text = "Modify (Magnify)"
content.text = ("What modifications can be made to improve performance?\n"
                "Example: Modifying the shape for an optimized sweet spot.\n"
                "Benefits: Enhanced hitting power and precision.")
# Download and add bat shape diagram image
download_image('https://static.wixstatic.com/media/2e2d5b_4c536fb28fef4663a43a105fefc213b7~mv2.jpg', 'bat_shape_modifications.jpg')
slide_6.shapes.add_picture('bat_shape_modifications.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 6: Put to Another Use
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]
title.text = "Put to Another Use"
content.text = ("What other uses can a cricket bat have?\n"
                "Example: Fitness training (weighted bat exercises).\n"
                "Benefits: Versatile usage increases market appeal.")
# Download and add fitness training image
download_image('https://image.shutterstock.com/image-photo/cricket-bat-isolated-white-background-260nw-1888110708.jpg', 'fitness_training.jpg')
slide_7.shapes.add_picture('fitness_training.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 7: Eliminate
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]
title.text = "Eliminate"
content.text = ("What aspects can be eliminated to improve usability?\n"
                "Example: Removing excess weight.\n"
                "Benefits: Easier handling without compromising strength.")
# Download and add visual of lightweight bats
download_image('https://www.woodbat.com/wp-content/uploads/2021/06/Untitled-design-11-300x300.png', 'lightweight_bat.jpg')
slide_8.shapes.add_picture('lightweight_bat.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 8: Rearrange (Reverse)
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.placeholders[1]
title.text = "Rearrange (Reverse)"
content.text = ("How can we rearrange the design?\n"
                "Example: Reversing the handle design for a better grip.\n"
                "Benefits: Improved ergonomics and comfort.")
# Download and add new handle design image
download_image('https://cdn.pixabay.com/photo/2020/01/12/15/26/cricket-4754827_1280.jpg', 'new_handle_design.jpg')
slide_9.shapes.add_picture('new_handle_design.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 9: Conclusion
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]
title.text = "Conclusion"
content.text = ("Recap of SCAMPER application to cricket bats.\n"
                "Encouragement to explore innovative ideas in sports equipment design.")
# Download and add inspirational quote image
download_image('https://www.wittonpark.co.uk/wp-content/uploads/2020/01/Inspiration.jpg', 'inspirational_quote.jpg')
slide_10.shapes.add_picture('inspirational_quote.jpg', Inches(1), Inches(2), width=Inches(5))

# Slide 10: Questions and Discussion
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_11.shapes.title
content = slide_11.placeholders[1]
title.text = "Questions and Discussion"
content.text = ("Invite questions from the audience.\n"
                "Open the floor for discussion on ideas generated using SCAMPER.")
# Download and add discussion image
download_image('https://image.shutterstock.com/image-photo/business-people-discussing-project-ideas-260nw-1226481730.jpg', 'discussion.jpg')
slide_11.shapes.add_picture('discussion.jpg', Inches(1), Inches(2), width=Inches(5))

# Save the presentation
prs.save('scamper_cricket_bat_with_images.pptx')

print("Presentation created successfully!")
